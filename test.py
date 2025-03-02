from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.xmlchemy import OxmlElement

from functions import *

# Create a presentation
prs = Presentation()

slide_layout = prs.slide_layouts[6]  # Blank slide layout

slide1 = prs.slides.add_slide(slide_layout)
slide_width, slide_height = prs.slide_width, prs.slide_height
center_width, center_height = slide_width / 2, slide_height / 2

add_header_and_footer(prs, slide1, "")
add_textbox(slide1, Inches(0.5), Inches(2), Inches(3), Inches(0.5), "Issued On: 27 February 2025", 12)
add_image(slide1, Inches(3.5), Inches(2), Inches(6), "images.jpg")

slide2 = prs.slides.add_slide(slide_layout)

bullet_points = [

]

info_box(slide2, Inches(0.2), Inches(1), Inches(4), Inches(0.3), Inches(0.9), "Summary", bullet_points, 12)
add_header_and_footer(prs, slide2, "Summary")
table, _ = create_blank_table(slide2, 6, 4, Inches(0.2), Inches(2.6), Inches(3.1), Inches(0.1))

set_cell_colour(table, 1, 0, RGBColor(255, 0, 0))
set_cell_colour(table, 2, 0, RGBColor(255, 64, 0))
set_cell_colour(table, 3, 0, RGBColor(255, 128, 0))
set_cell_colour(table, 4, 0, RGBColor(255, 191, 0))
set_cell_colour(table, 5, 0, RGBColor(255, 255, 0))

table.columns[1].width = Inches(1.8)

populate_table(table, "p.csv")

set_table_font_size(table, 9)
add_line(prs, slide2, Inches(0.2), slide_width - Inches(0.2), Inches(2.4))

slide3 = prs.slides.add_slide(slide_layout)
add_header_and_footer(prs, slide3, "")

add_textbox(slide3, center_width + Inches(0.15), Inches(1.1), Inches(1.25), Inches(0.2), "Consequences", 11, bold = True)

table1, table_shape = create_blank_table(slide3, 6, 4, Inches(5.1), Inches(1.4), Inches(3.9), Inches(2.4))

set_cell_colour(table1, 1, 0, RGBColor(255, 0, 0))
set_cell_colour(table1, 2, 0, RGBColor(255, 64, 0))
set_cell_colour(table1, 3, 0, RGBColor(255, 128, 0))
set_cell_colour(table1, 4, 0, RGBColor(255, 191, 0))
set_cell_colour(table1, 5, 0, RGBColor(255, 255, 0))

populate_table(table1, "p.csv")
table1.columns[1].width = Inches(1.8)  # Set the second column width to 2.5 inches
set_table_font_size(table1, 7)

bullet_points = [

]

b_box, b_frame = add_bullet_points(slide3, bullet_points, Inches(5.1), table_shape.height + Inches(1.45), Inches(4.8), Inches(1.45), 8)

bullet_points = [

]

add_line(prs, slide3, Inches(5.1), Inches(9.8), Inches(5.2))

info_box(slide3, Inches(5.1), Inches(5.3), Inches(4.7), Inches(0.2), Inches(0.6), "Scenario", bullet_points, 8)

bullet_points = [

]

info_box(slide3, Inches(5.1), Inches(6.2), Inches(4.7), Inches(0.2), Inches(0.5), "Assumptions", bullet_points, 8)

prs.save("presentation_with_header.pptx")
