from typing import List

from pptx.slide import Slide
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame
from pptx.dml.color import RGBColor

def add_textbox(slide: Slide, left: Inches, top: Inches, width: Inches, height: Inches, text: str, font_size: int, center: bool = False, bold: bool = False, word_wrap: bool = False) -> tuple[BaseShape, TextFrame]:
    """Adds a textbox to a slide with specified text formatting options."""

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = word_wrap

    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(font_size)
    
    # Set bold if required
    if bold:
        run.font.bold = True

    if center:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
    
    return text_box, text_frame

def add_bullet_points(slide: Slide, bullet_points: List[str], left: Inches, top: Inches, width: Inches, height: Inches, font_size: int) -> tuple[BaseShape, TextFrame]:
    """Adds a bulleted list to a slide and returns the shape and text frame."""

    paragraph_box = slide.shapes.add_textbox(left, top, width, height)
    paragraph_frame = paragraph_box.text_frame

    paragraph_frame.word_wrap = True 

    for point in bullet_points:
        if paragraph_frame.paragraphs[-1].text:
            bullet_p = paragraph_frame.add_paragraph()
        else:
            bullet_p = paragraph_frame.paragraphs[-1]
        bullet_p.text = f"{point}"
        bullet_p.font.size = Pt(font_size)  # Set font size
    
    return paragraph_box, paragraph_frame

def add_info_box(slide: Slide, left: Inches, top: Inches, width: Inches, header_text: str, paragraph_text: List[str], font_size: int, padding: Inches = Inches(0)) -> tuple[BaseShape, TextFrame, BaseShape, TextFrame]:
    """Adds an information box to a slide consisting of a header and a paragraph of bullet points."""

    header_height = Inches(0.25)
    header_box = slide.shapes.add_textbox(left, top, width, header_height) 
    header_frame = header_box.text_frame
    header_frame.text = header_text

    header_p = header_frame.paragraphs[0]
    header_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    header_p.font.size = Pt(font_size + 1)

    # Set the background color 
    header_fill = header_box.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(117, 128, 139)  
    header_height = header_box.height

    # Add the paragraph text box below the header

    paragraph_height = Inches(estimate_bullet_point_textbox_height(paragraph_text, font_size, int(width / 914400))) + Inches(0.05) + padding
    paragraph_box = slide.shapes.add_textbox(left, top + header_height, width, paragraph_height)
    paragraph_frame = paragraph_box.text_frame

    paragraph_frame.word_wrap = True  # Ensures text wraps inside the box

    for point in paragraph_text:
        if paragraph_frame.paragraphs[-1].text:
            bullet_p = paragraph_frame.add_paragraph()
        else:
            bullet_p = paragraph_frame.paragraphs[-1]
        bullet_p.text = point
        bullet_p.font.size = Pt(font_size)  #

    # Format the paragraph background color
    paragraph_fill = paragraph_box.fill
    paragraph_fill.solid()
    paragraph_fill.fore_color.rgb = RGBColor(233, 233, 233)  # Grey

    return header_box, header_frame, paragraph_box, paragraph_frame

def estimate_textbox_height(text: str, font_size_pt: int, max_width_inch: int) -> int:
    """Estimate the height of a text box in PowerPoint for Calibri (Body) font."""

    avg_char_width_inch = (font_size_pt * 0.5) / 72  # Approx 50% of font size, converted to inches
    line_height_inch = (font_size_pt * 1.25) / 72  # Line height ~1.25× font size, converted to inches

    # Estimate the number of lines based on word wrapping
    words = text.split()
    lines = []
    current_line = ""

    for word in words:
        test_line = current_line + " " + word if current_line else word
        if len(test_line) * avg_char_width_inch > max_width_inch:
            lines.append(current_line)
            current_line = word
        else:
            current_line = test_line

    if current_line:
        lines.append(current_line)  # Add last line

    total_height_inch = len(lines) * line_height_inch
    return total_height_inch

def estimate_bullet_point_textbox_height(bullet_points: List[str], font_size_pt: int, max_width_inch: int) -> int:
    """Estimates the height of a text box containing a list of bullet points based on font size and width constraints."""

    height = 0

    for point in bullet_points:
        height += estimate_textbox_height(point, font_size_pt, max_width_inch)
    return height