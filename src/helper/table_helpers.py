import pandas as pd

from typing import Any, List

from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.util import Pt, Inches

from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame
from pptx.table import Table

from pptx.oxml.xmlchemy import OxmlElement
from pptx.table import _Cell

from .text_helpers import add_textbox

def calculate_table_height(table: Table) -> Inches:
    """Calculates the total height of a table based on its number of rows and font size."""

    font_size = int(table.cell(0, 0).text_frame.paragraphs[0].font.size.pt)
    rows = len(table.rows)
    
    if font_size == 7:
        return Inches(rows * 0.219)

def add_table(slide: Slide, title: str, csv_data: str, cell_colours: List[RGBColor], top: Inches, left: Inches, table_width: Inches) -> tuple[Table, BaseShape]:
    """
    Adds a table with a title to a slide.

    The function creates and populates a table, sets column width, applies cell colors,
    adjusts font size, and calculates the table height. Returns the table and its shape.
    """

    add_textbox(slide, left, top, Inches(0.5), Inches(0.15), title, 9, bold = True)
    table, shape = create_and_populate_table(slide, left + Inches(0.05), top + Inches(0.22), table_width, Inches(0.1), csv_data)
    table.columns[1].width = Inches(1.8)

    set_cell_colours(table, cell_colours)
    set_table_font_size(table, 7)
    table.height = calculate_table_height(table)

    return table, shape 

def create_and_populate_table(slide: Slide, x: Inches, y: Inches, cx: Inches, cy: Inches, csv: str) -> tuple[Table, BaseShape]:
    """
    Creates a blank table on a slide and populates it with data from a CSV file.

    The table includes an extra row for headers and an empty first column.
    Returns the table and its shape.
    """
    df = pd.read_csv(csv)
    rows = df.shape[0] + 1 # Plus 1 for the headers
    cols = df.shape[1] + 1 # Plus 1 for empty first colunmn

    table, table_shape = create_blank_table(slide, rows, cols, x, y, cx, cy)
    populate_table(table, csv)

    return table, table_shape

def populate_table(table: Table, csv_path: str) -> None:
    """
    Populates a table with data from a CSV file.

    The first row is filled with column headers, and subsequent rows are populated with CSV data.
    """

    # Read the CSV file into a DataFrame
    df = pd.read_csv(csv_path)
    
    # Populate the first row with column headers
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx + 1).text = col_name  # Adjust the index if necessary (depends on table implementation)
    
    # Iterate through each row and column index to fill the table with data
    for row_idx, row in enumerate(df.itertuples(index=False), start=1):  # Start from 1 to skip the header row
        for col_idx, cell in enumerate(row):
            table.cell(row_idx, col_idx + 1).text = str(cell)  # Populate the table with the cell data

def create_blank_table(slide: Slide, rows: int, columns: int, x: Inches, y: Inches, cx: Inches, cy: Inches) -> tuple[Table, BaseShape]:
    """
    Creates a blank table on a slide with a specified number of rows and columns.

    The table cells are initialised with borders, a white background, and a default font size.
    """

    shape = slide.shapes.add_table(rows, columns, x, y, cx, cy)
    table = shape.table

    for r in range(rows):
        for c in range(columns):
                cell = table.cell(r, c)
                _set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    set_table_font_size(table, 10)
    return table, shape

def set_table_font_size(table: Table, size: int) -> None:
    """
    Sets the font size and colour for all text in a table.

    Applies the specified font size and ensures text colour is black.
    """

    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(size)
                for run in paragraph.runs:
                    run.font.size = Pt(size) 
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Black colour

def _SubElement(parent: OxmlElement, tagname: str, **kwargs: Any) -> OxmlElement:
    """Creates and appends a sub-element to a parent element."""

    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell: _Cell, border_color: str = "000000", border_width: str = "12700") -> None:
    """Sets the border for a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = _SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = _SubElement(ln, 'a:solidFill')
        srgbClr = _SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = _SubElement(ln, 'a:prstDash', val='solid')
        round_ = _SubElement(ln, 'a:round')
        headEnd = _SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = _SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def set_cell_colour(table: Table, row: int, column: int, colour: List[RGBColor]):
    """
    Sets the background colour of a specific cell in a table.

    The specified colour is applied to the cell at the given row and column indices.
    """
    specific_cell = table.cell(row, column) 
    specific_cell.fill.solid()  # Apply solid fill
    specific_cell.fill.fore_color.rgb = colour 

def set_cell_colours(table: Table, colours: List[RGBColor]):
    """
    Sets the background colour for the first column of each row in a table.

    The colors are applied to the cells in the first column based on the provided list of colors.
    """
    for i in range(len(colours)):
        if i + 2 > len(table.rows):
            break 
        set_cell_colour(table, i + 1, 0, colours[i])