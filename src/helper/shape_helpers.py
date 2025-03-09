from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.slide import Slide
from pptx.shapes.shapetree import SlideShapes

def add_line(slide: Slide, start_x: Inches, end_x: Inches, y: Inches) -> 'SlideShapes._Shape':
    """
    Adds a horizontal line to a slide.

    The line is a straight connector with a black color and a thin width.
    """
    # Add a straight connector (line) across the slide
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, y, end_x, y)

    # Change the color of the connector line to black
    line.line.color.rgb = RGBColor(0, 0, 0) 

    # Set the thickness of the line
    line.line.width = Inches(0.01) 

    return line