from pptx.util import Inches
from pptx.slide import Slide
from pptx.shapes.shapetree import SlideShapes

def add_image(slide: Slide, left: Inches, top: Inches, width: Inches, path: str, height: Inches = Inches(0)) -> 'SlideShapes._Shape':
    """
    Adds an image to a slide at the specified position and size.

    If a height is provided, both width and height are used to scale the image. 
    Otherwise, only the width is used, and the height is scaled automatically.
    """
    
    if height != Inches(0):
        return slide.shapes.add_picture(path, left, top, width, height) 
    return slide.shapes.add_picture(path, left, top, width) 