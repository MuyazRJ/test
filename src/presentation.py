from typing import Type

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.slide import Slide

from .frame.base import BaseSlide
from .frame.title import TitleSlide
from .frame.new import NewSlide
from.frame.summary import SummarySlide

from .helper.text_helpers import add_textbox
from .helper.shape_helpers import add_line

class BriefingPack():
    """
    Class to manage a PowerPoint presentation for a briefing pack.

    This class facilitates the creation and management of slides in a PowerPoint presentation.
    It handles adding slide frames (e.g., title, table, image), numbering slides, and adding common templates (header, footer) to each slide.
    """
    
    def __init__(self, reference_number: str, classification: str, code_version: str, job_id: str) -> None:
        self.prs = Presentation()
        self.reference_number = reference_number
        self.classification = classification
        self.code_version = code_version
        self.job_id = job_id
    
    def add_frame(self, frame: Type['BaseSlide']) -> None:
        """
        Adds a slide frame to the briefing pack by invoking the `add_slide` method of the given slide frame class
        """

        slide = frame.add_slide(self)
        self.add_slide_template(slide, frame.slide_title)

    def number_slides(self):
        """
        Adds slide numbers to each slide in the presentation, formatted as "X/Y" (e.g., "1/5").
        The slide number appears in the bottom-right corner of each slide.
        """

        slide_total = len(self.prs.slides)

        for index, slide in enumerate(self.prs.slides):
            # Define the position where the slide number will appear
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
            slide_number = index + 1
            
            # Adjust position for the slide number (bottom-right corner)
            left = slide_width - Inches(0.35)
            top = slide_height - Inches(0.2)

            # Add a text box with the slide number
            add_textbox(slide, left, top, Inches(3), Inches(0.2), f"{slide_number}/{slide_total}", 7)
    
    def add_slide_template(self, slide: Slide, slide_title: str) -> None:
        """
        Adds a header and footer to the given slide, including reference number, classification, dividing line,
        code version, job ID, and slide title.
        """
       
        slide_width = self.prs.slide_width

        add_textbox(slide, Inches(0.1), Inches(0), Inches(1), Inches(1), f"Reference Number: {self.reference_number}", 7)

        # Add classification to the top middle
        center_x = (slide_width - Inches(1.5)) / 2  # Center the text box horizontally
        add_textbox(slide, center_x, Inches(0), Inches(1.5), Inches(0.3), self.classification, 7, center = True, bold =True)

        # Calculate the start and end points for the line
        line_start_x = Inches(0.2)  
        line_end_x = slide_width - Inches(0.2)  
        line_y = Inches(0.7) 
        add_line(slide, line_start_x, line_end_x, line_y)

        bottom_y = self.prs.slide_height - Inches(0.2)

        # Add text to the bottom left 
        add_textbox(slide, Inches(0.1), bottom_y, Inches(1), Inches(1), f"Code version: {self.code_version}", 7)
        add_textbox(slide, Inches(1), bottom_y, Inches(1), Inches(1), f"Job ID: {self.job_id}", 7)

        # Add classification to the bottom center 
        bottom_center_x = (slide_width - Inches(1.5)) / 2  # Center the text box horizontally
        add_textbox(slide, bottom_center_x, bottom_y, Inches(1.5), Inches(0.3), self.classification, 7, center = True, bold = True)

        # Add slide title above the black line
        add_textbox(slide, Inches(0.11), Inches(0.36), Inches(1), Inches(0.4), slide_title, 18)

    def save(self, filename: str) -> None:
        """Save the presentation to a file."""
        self.number_slides()
        self.prs.save(filename)


