from .base import BaseSlide

from pptx.util import Inches
from pptx.slide import Slide

from ..helper.text_helpers import add_textbox
from ..helper.image_helpers import add_image

class TitleSlide(BaseSlide):
    """
    Represents a title slide with an issue date and an image.

    This class handles the creation of a title slide that includes:
    - A textbox displaying the issue date.
    - An image positioned at a specific location with a defined width.
    """

    def __init__(self, slide_title: str, issue_date: str, image_path: str, image_x: Inches = Inches(3.5), image_width: Inches = Inches(6)) -> None:
        super().__init__(slide_title)
        self.issue_date, self.image_path = issue_date, image_path
        self.image_x, self.image_width = image_x, image_width
    
    def add_slide(self, briefing_pack) -> Slide:
        """
        Adds a title slide to the briefing pack, including the issue date and image.
        """

        self.prs = briefing_pack.prs
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Create an empty slide

        add_textbox(self.slide, Inches(0.5), Inches(2), Inches(3), Inches(0.5), f"Issued On: {self.issue_date}", 12)
        add_image(self.slide, self.image_x, Inches(2), self.image_width, self.image_path)

        return self.slide