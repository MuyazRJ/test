from .base import BaseSlide
from typing import List

from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.slide import Slide

from ..helper.text_helpers import add_info_box, add_bullet_points, estimate_bullet_point_textbox_height
from ..helper.shape_helpers import add_line
from ..helper.json_helpers import load_bullet_points
from ..helper.image_helpers import add_image
from ..helper.table_helpers import add_table

class NewSlide(BaseSlide):
    """
    Represents a slide with a table, bullet points, and an image.

    This class handles the creation of a slide that includes:
    - A table populated with data from a CSV file.
    - Bullet points for the table, scenario, and assumptions.
    - An image inserted into the slide.
    """

    def __init__(self, slide_title: str, table_header: str, table_csv: str, table_cell_colours: List[RGBColor], table_bullet_points_key: str, scenario_bullet_points_key: str, assumptions_bullet_points_key: str, image_path: str) -> None:
        super().__init__(slide_title)

        self.table_header = table_header
        self.table_csv, self.table_cell_colours = table_csv, table_cell_colours

        self.table_bullet_points = load_bullet_points("input/bullet_points.json", table_bullet_points_key)
        self.scenario_bullet_points= load_bullet_points("input/bullet_points.json", scenario_bullet_points_key)
        self.assumptions_bullet_points = load_bullet_points("input/bullet_points.json", assumptions_bullet_points_key)

        self.image_path = image_path
    
    def add_slide(self, briefing_pack) -> Slide:
        """Adds a slide to the briefing pack with a table, bullet points, and an image."""
        
        self.prs = briefing_pack.prs
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Create an empty slide

        content_left = Inches(5.02)

        table, table_shape = add_table(self.slide, self.table_header, self.table_csv, self.table_cell_colours, Inches(1), content_left, Inches(3.9))

        table_text_box, table_text_frame = add_bullet_points(self.slide, self.table_bullet_points, content_left, table_shape.top + table.height + Inches(0.05), Inches(4.78), Inches(estimate_bullet_point_textbox_height(self.table_bullet_points, 8, 4.78)), 8)
        line = add_line(self.slide, content_left + Inches(0.095), Inches(9.78), table_text_box.top + table_text_box.height + Inches(0.07))

        _, _, sc_p_box, _ = add_info_box(self.slide, content_left + Inches(0.05), line.top + Inches(0.07), Inches(4.72), "Scenario", self.scenario_bullet_points, 8)
        add_info_box(self.slide, content_left + Inches(0.05), sc_p_box.top + sc_p_box.height + Inches(0.1), Inches(4.72), "Assumptions", self.assumptions_bullet_points, 8)

        add_image(self.slide, Inches(0.2), Inches(1.1), Inches(4.7), self.image_path, Inches(5.9))

        return self.slide