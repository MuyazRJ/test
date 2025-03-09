from .base import BaseSlide
from typing import List

from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.slide import Slide

from ..helper.text_helpers import add_textbox, add_info_box
from ..helper.shape_helpers import add_line
from ..helper.json_helpers import load_bullet_points
from ..helper.image_helpers import add_image
from ..helper.table_helpers import add_table

class SummarySlide(BaseSlide):
    """
    Represents a summary slide with two tables, bullet points, comments, and optional images.

    This class handles the creation of a summary slide, which includes:
    - Bullet points in an information box.
    - One or two tables populated with data from CSV files.
    - Optional images added beside the tables.
    - A comment section.
    """

    def __init__(self, slide_title: str, bullet_point_key: str, comments: str, table1_title: str, table1_csv: str, table1_cell_colours: List[RGBColor], table2_title: str, table2_csv: str, table2_cell_colours: List[RGBColor], image1_path: str, image2_path: str) -> None:
        super().__init__(slide_title)

        self.comments = comments
        self.bullet_points = load_bullet_points("input/bullet_points.json", bullet_point_key)

        self.table1_csv, self.table2_csv = table1_csv, table2_csv
        self.table1_cell_colours, self.table2_cell_colours = table1_cell_colours, table2_cell_colours
        self.table1_title, self.table2_title = table1_title, table2_title

        # Placeholders for tables and their shapes
        self.table1, self.table1_shape = None, None
        self.table2, self.table2_shape = None, None
        
        self.image1_path, self.image2_path = image1_path, image2_path

    def add_slide(self, briefing_pack) -> Slide:
        """Adds a summary slide to the briefing pack, including bullet points, tables, comments, and images."""

        self.prs = briefing_pack.prs
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Create an empty slide

        # Add scenario info box and dividing line
        add_info_box(self.slide, Inches(0.2), Inches(1), Inches(5), "Scenario", self.bullet_points, 9, padding = Inches(0.8))
        add_line(self.slide, Inches(0.2), self.prs.slide_width - Inches(0.2), Inches(2.85))

        # Comment area
        add_textbox(self.slide, Inches(5.2), Inches(0.95), Inches(4.6), Inches(1.7), f"Comments: {self.comments}", 9, word_wrap = True)

        # Add first table
        self.table1, self.table1_shape = add_table(self.slide, self.table1_title, self.table1_csv, self.table1_cell_colours, Inches(3), Inches(0.15), Inches(3))

        # Add second table if applicable
        if self.table2_csv:
            top = self.table1.height + self.table1_shape.top + Inches(0.2)
            self.table2, self.table2_shape = add_table(self.slide, self.table2_title, self.table2_csv, self.table2_cell_colours, top, Inches(0.15), Inches(3))

        self._add_images()
        return self.slide

    def _add_images(self) -> None:
        """
        Adds one or two images beside the tables in the slide.

        If both images are provided, they will be displayed side by side. Otherwise, a single image is added.
        """
         
        if not self.table1_shape:
            return  # No table available to align images with

        image_area = self.prs.slide_width - (self.table1_shape.width + self.table1_shape.left) - Inches(0.5)
        image_left = self.table1_shape.width + self.table1_shape.left + Inches(0.2)
        image_height = (self.table2_shape.top + self.table2.height - self.table1_shape.top + Inches(0.7)) if self.table2_shape else Inches(2)

        if self.image2_path:
            # Add two images side by side
            half_width = image_area / 2
            add_image(self.slide, image_left, Inches(3.22), half_width, self.image1_path, image_height)
            add_image(self.slide, image_left + half_width + Inches(0.1), Inches(3.22), half_width, self.image2_path, image_height)
        else:
            # Add a single image
            add_image(self.slide, image_left, Inches(3.22), image_area, self.image1_path, image_height)