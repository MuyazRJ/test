from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import MSO_AUTO_SIZE

import pandas as pd

def add_textbox(slide, left, top, width, height, text, font_size, center = False, bold = False):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text

    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(font_size)
    
    # Set bold if required
    if bold:
        run.font.bold = True

    if center:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

def add_line(prs, slide, start_x, end_x, y):
    # Add a straight connector (line) across the slide
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, y, end_x, y)

    # Change the color of the connector line to black
    line1.line.color.rgb = RGBColor(0, 0, 0)  # Black color

    # Set the thickness of the line
    line1.line.width = Inches(0.01) 

def add_header_and_footer(prs, slide, title = ""):
    """Function to add a header and footer to the slide."""
    slide_width = prs.slide_width

    add_textbox(slide, Inches(0.1), Inches(0), Inches(1), Inches(1), "Reference Number: ", 7)

    # Add text to the top middle (centered horizontally, 0.5 inch from the top)
    center_x = (slide_width - Inches(1.5)) / 2  # Center the text box horizontally
    add_textbox(slide, center_x, Inches(0), Inches(1.5), Inches(0.3), "OFFICIAL-SENSITIVE", 7, True, True)

    # Calculate the start and end points for the line
    line_start_x = Inches(0.2)  # 0.5 inch from the left
    line_end_x = slide_width - Inches(0.2)  # 0.5 inch from the right
    line_y = Inches(0.7)  # Set the Y position of the line (0.7 inch from the top)
    add_line(prs, slide, line_start_x, line_end_x, line_y)

    # Add text to the bottom left (0.5 inch from the left and 0.5 inch from the bottom)
    bottom_y = prs.slide_height - Inches(0.2)

    # Add text to the top left (0.5 inch from the left and 0.5 inch from the top)
    add_textbox(slide, Inches(0.1), bottom_y, Inches(1), Inches(1), "Code version: 5.0.0", 7)
    add_textbox(slide, Inches(1), bottom_y, Inches(1), Inches(1), "Job ID: strike-2025-02-19.13-24-41", 7)

    # Add text to the bottom center (centered horizontally, 0.5 inch from the bottom)
    bottom_center_x = (slide_width - Inches(1.5)) / 2  # Center the text box horizontally
    add_textbox(slide, bottom_center_x, bottom_y, Inches(1.5), Inches(0.3), "OFFICIAL-SENSITIVE", 7, True, True)

    if title:
        add_textbox(slide, Inches(0.11), Inches(0.36), Inches(1), Inches(0.4), title, 18)

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
            
def create_blank_table(slide, rows, columns, x, y, cx, cy):
    shape = slide.shapes.add_table(rows, columns, x, y, cx, cy)
    table = shape.table

    for r in range(rows):
        for c in range(columns):
                cell = table.cell(r, c)
                _set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    set_table_font_size(table, 10)
    return table, shape

def set_cell_colour(table, row, column, colour):
    specific_cell = table.cell(row, column)  # Access the cell in the second row, second column
    specific_cell.fill.solid()  # Apply solid fill
    specific_cell.fill.fore_color.rgb = colour  # Set the background color to r

def set_table_font_size(table, size):
     for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(size)  # Set font size to 14pt
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Red color (RGB)

def info_box(slide, left, top, width, header_height, paragraph_height, header_text, paragraph_text, font_size):
    header_box = slide.shapes.add_textbox(left, top, width, header_height)  # Initial small height
    header_frame = header_box.text_frame
    header_frame.text = header_text

    header_p = header_frame.paragraphs[0]
    header_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    header_p.font.size = Pt(font_size + 1)

    # Set the background color (Blue)
    header_fill = header_box.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(117, 128, 139)  # Blue

    # Auto-size the header box
    header_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Resize to fit text
    header_height = header_box.height

    # Add the paragraph text box below the header
    paragraph_box = slide.shapes.add_textbox(left, top + header_height, width, paragraph_height)
    paragraph_frame = paragraph_box.text_frame

    paragraph_frame.word_wrap = True  # Ensures text wraps inside the box
    paragraph_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Prevents auto-expansion

    for point in paragraph_text:
        if paragraph_frame.paragraphs[-1].text:
            bullet_p = paragraph_frame.add_paragraph()
        else:
            bullet_p = paragraph_frame.paragraphs[-1]
        bullet_p.text =  f"• {point}"
        bullet_p.font.size = Pt(font_size)  # Set font size

    # Format the paragraph background color (Grey)
    paragraph_fill = paragraph_box.fill
    paragraph_fill.solid()
    paragraph_fill.fore_color.rgb = RGBColor(233, 233, 233)  # Grey

def add_image(slide, left, top, width, path):
    return slide.shapes.add_picture(path, left, top, width) 

def add_bullet_points(slide, bullet_points, left, top, width, height, font_size):
    paragraph_box = slide.shapes.add_textbox(left, top, width, height)
    paragraph_frame = paragraph_box.text_frame

    paragraph_frame.word_wrap = True  # Ensures text wraps inside the box
    paragraph_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Prevents auto-expansion

    for point in bullet_points:
        if paragraph_frame.paragraphs[-1].text:
            bullet_p = paragraph_frame.add_paragraph()
        else:
            bullet_p = paragraph_frame.paragraphs[-1]
        bullet_p.text = f"• {point}"
        bullet_p.font.size = Pt(font_size)  # Set font size
    
    return paragraph_box, paragraph_frame

def populate_table(table, csv_path):
    # Read the CSV file into a DataFrame
    df = pd.read_csv(csv_path)
    
    # Populate the first row with column headers
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx + 1).text = col_name  # Adjust the index if necessary (depends on table implementation)
    
    # Iterate through each row and column index to fill the table with data
    for row_idx, row in enumerate(df.itertuples(index=False), start=1):  # Start from 1 to skip the header row
        for col_idx, cell in enumerate(row):
            table.cell(row_idx, col_idx + 1).text = str(cell)  # Populate the table with the cell data

